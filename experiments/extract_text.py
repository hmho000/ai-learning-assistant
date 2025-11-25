"""
extract_text.py
----------------

通用文本抽取层，支持：
- PDF → 文本（优先 PyMuPDF，若文本不足则自动 OCR）
- 图片 PDF / 扫描 PDF → OCR
- 图片文件（png/jpg/jpeg/webp）
- Word (docx)
- PPT (pptx)

OCR 优先调用 SiliconFlow 多模态 API，若不可用则回退到 PaddleOCR。
"""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path
from typing import Callable, List, Optional

import fitz  # PyMuPDF
import requests

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    from paddleocr import PaddleOCR
except ImportError:  # pragma: no cover
    PaddleOCR = None


MIN_TEXT_THRESHOLD = 50
SUPPORTED_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
SILICONFLOW_API_URL = "https://api.siliconflow.cn/v1/images/ocr"

_paddle_ocr_client: Optional[PaddleOCR] = None


def extract_text_from_file(path: str) -> str:
    """
    自动根据文件类型（pdf / docx / pptx / 图片）提取可读文本。
    """
    filepath = Path(path)
    if not filepath.exists():
        raise FileNotFoundError(f"未找到文件：{filepath}")

    suffix = filepath.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(filepath)
    if suffix == ".docx":
        return _extract_docx(filepath)
    if suffix == ".pptx":
        return _extract_pptx(filepath)
    if suffix in SUPPORTED_IMAGE_SUFFIXES:
        return _extract_image(filepath)

    raise ValueError(f"暂不支持的文件类型：{suffix}")


# -----------------------------------------------------------------------------
# PDF 处理
# -----------------------------------------------------------------------------

def _extract_pdf(pdf_path: Path) -> str:
    text = _extract_pdf_with_pymupdf(pdf_path)
    if len(text.strip()) >= MIN_TEXT_THRESHOLD:
        return text
    return _extract_pdf_via_ocr(pdf_path)


def _extract_pdf_with_pymupdf(pdf_path: Path) -> str:
    doc = fitz.open(pdf_path)
    try:
        chunks: List[str] = []
        for page in doc:
            chunks.append(page.get_text("text"))
        return "\n".join(chunks).strip()
    finally:
        doc.close()


def _extract_pdf_via_ocr(pdf_path: Path) -> str:
    doc = fitz.open(pdf_path)
    try:
        page_results: List[str] = []
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            page_results.append(_run_ocr(img_bytes))
        return "\n".join(page_results).strip()
    finally:
        doc.close()


# -----------------------------------------------------------------------------
# DOCX 处理
# -----------------------------------------------------------------------------

def _ensure_docx():
    if DocxDocument is None:
        raise RuntimeError("缺少 python-docx 依赖，请先安装：pip install python-docx")


def _extract_docx(docx_path: Path) -> str:
    _ensure_docx()
    doc = DocxDocument(docx_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


# -----------------------------------------------------------------------------
# PPTX 处理
# -----------------------------------------------------------------------------

def _ensure_pptx():
    if Presentation is None:
        raise RuntimeError("缺少 python-pptx 依赖，请先安装：pip install python-pptx")


def _extract_pptx(pptx_path: Path) -> str:
    _ensure_pptx()
    prs = Presentation(pptx_path)
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
    return "\n".join(texts).strip()


# -----------------------------------------------------------------------------
# 图片 / OCR
# -----------------------------------------------------------------------------

def _extract_image(image_path: Path) -> str:
    with image_path.open("rb") as f:
        return _run_ocr(f.read())


def _run_ocr(image_bytes: bytes) -> str:
    ocr_functions: List[Callable[[bytes], Optional[str]]] = [
        _ocr_via_siliconflow,
        _ocr_via_paddleocr,
    ]
    for func in ocr_functions:
        try:
            result = func(image_bytes)
            if result:
                return result
        except Exception as exc:  # pragma: no cover - 记录错误但继续
            print(f"[OCR] {func.__name__} 失败：{exc}")
            continue
    raise RuntimeError("所有 OCR 后端均不可用，请检查配置。")


def _ocr_via_siliconflow(image_bytes: bytes) -> Optional[str]:
    api_key = os.getenv("SILICONFLOW_API_KEY")
    if not api_key:
        return None
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "image_base64": base64.b64encode(image_bytes).decode("utf-8"),
        "language": "auto",
    }
    resp = requests.post(SILICONFLOW_API_URL, json=payload, headers=headers, timeout=120)
    resp.raise_for_status()
    data = resp.json()
    return data.get("text")


def _ocr_via_paddleocr(image_bytes: bytes) -> Optional[str]:
    global _paddle_ocr_client
    if PaddleOCR is None:
        return None
    if _paddle_ocr_client is None:
        _paddle_ocr_client = PaddleOCR(use_angle_cls=True, lang="ch")
    image_stream = io.BytesIO(image_bytes)
    result = _paddle_ocr_client.ocr(image_stream, cls=True)
    if not result:
        return ""
    lines = []
    for line in result:
        if not line:
            continue
        text = line[0][1][0]
        lines.append(text)
    return "\n".join(lines).strip()


if __name__ == "__main__":  # 简易测试
    import argparse

    parser = argparse.ArgumentParser(description="通用文本抽取测试脚本")
    parser.add_argument("path", type=str, help="待抽取的文件路径")
    args = parser.parse_args()
    print(extract_text_from_file(args.path))

